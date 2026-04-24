# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(r"C:\Users\12700\email_llm_assistant\毕业设计预答辩_高煜同.pptx")
FALLBACK_OUT = Path(r"C:\Users\12700\email_llm_assistant\毕业设计预答辩_高煜同_新生成.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(29, 78, 216)
BLUE_DARK = RGBColor(30, 64, 175)
BLUE_LIGHT = RGBColor(239, 246, 255)
GRAY_BG = RGBColor(243, 244, 246)
GRAY_LINE = RGBColor(229, 231, 235)
GRAY_TEXT = RGBColor(75, 85, 99)
LIGHT_TEXT = RGBColor(156, 163, 175)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(17, 24, 39)
GREEN = RGBColor(22, 163, 74)
YELLOW = RGBColor(202, 138, 4)
PURPLE = RGBColor(147, 51, 234)


def fontify(run, size=18, bold=False, color=BLACK, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_background(slide, color=RGBColor(245, 247, 250)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, fill=WHITE, line=GRAY_LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    fontify(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, size=18, box=False):
    if box:
        add_card(slide, left, top, width, height)
        left += Inches(0.2)
        top += Inches(0.2)
        width -= Inches(0.4)
        height -= Inches(0.35)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "• " + item
        p.space_after = Pt(8)
        for run in p.runs:
            fontify(run, size=size, color=GRAY_TEXT)
    return tx


def add_footer(slide):
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(249, 250, 251)
    footer.line.color.rgb = GRAY_LINE
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.45), Inches(5.8), Inches(0.2), "2026 本科毕业设计预答辩 - 郑州西亚斯学院", size=10, color=LIGHT_TEXT)
    add_text(slide, Inches(8.0), SLIDE_H - Inches(0.45), Inches(4.8), Inches(0.2), "数字技术产业学院 | 人工智能专业", size=10, color=LIGHT_TEXT, align=PP_ALIGN.RIGHT)


def add_slide_number(slide, current, total):
    add_text(slide, Inches(11.9), Inches(0.25), Inches(0.9), Inches(0.2), f"{current} / {total}", size=10, color=LIGHT_TEXT, align=PP_ALIGN.RIGHT)


def add_icon_circle(slide, left, top, diameter, fill, text):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_text(slide, left, top + Inches(0.02), diameter, Inches(0.18), text, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def title_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_card(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(6.1))
    add_icon_circle(slide, Inches(5.95), Inches(1.0), Inches(1.1), BLUE, "PPT")
    add_text(slide, Inches(1.0), Inches(2.05), Inches(11.3), Inches(0.9),
             "基于 Llama 3.2 1B 的本地智能邮件分类与回复助手",
             size=28, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.8), Inches(2.95), Inches(5.7), Inches(0.45),
             "本科毕业设计预答辩汇报", size=18, bold=True, color=BLUE)
    add_text(slide, Inches(4.0), Inches(4.2), Inches(2.0), Inches(0.25), "汇报人：高煜同", size=13, color=GRAY_TEXT)
    add_text(slide, Inches(7.0), Inches(4.2), Inches(2.0), Inches(0.25), "专  业：人工智能", size=13, color=GRAY_TEXT)
    add_text(slide, Inches(4.0), Inches(4.65), Inches(2.7), Inches(0.25), "学  院：数字技术产业学院", size=13, color=GRAY_TEXT)
    add_text(slide, Inches(7.0), Inches(4.65), Inches(2.0), Inches(0.25), "指导老师：高媛", size=13, color=GRAY_TEXT)
    add_slide_number(slide, idx, total)
    add_footer(slide)


def list_slide(prs, idx, total, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_card(slide, Inches(0.45), Inches(0.45), Inches(12.45), Inches(6.25))
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(1.05), Inches(0.08), Inches(0.7)).fill.solid()
    add_text(slide, Inches(1.1), Inches(1.0), Inches(6.0), Inches(0.6), title, size=24, bold=True, color=BLUE_DARK)
    for i, item in enumerate(items, start=1):
        y = Inches(1.9) + Inches((i - 1) * 0.58)
        add_icon_circle(slide, Inches(1.05), y + Inches(0.02), Inches(0.34), BLUE_DARK, str(i))
        add_card(slide, Inches(1.55), y, Inches(10.6), Inches(0.42), fill=GRAY_BG, line=GRAY_BG)
        add_text(slide, Inches(1.8), y + Inches(0.03), Inches(10.0), Inches(0.28), item, size=18, color=GRAY_TEXT)
    add_slide_number(slide, idx, total)
    add_footer(slide)


def split_slide(prs, idx, total, title, left_items, right_builder):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_card(slide, Inches(0.45), Inches(0.45), Inches(12.45), Inches(6.25))
    add_text(slide, Inches(0.95), Inches(0.95), Inches(6.2), Inches(0.5), title, size=24, bold=True, color=BLUE_DARK)
    add_card(slide, Inches(0.95), Inches(1.8), Inches(5.35), Inches(4.35))
    y = Inches(2.15)
    for item in left_items:
        add_text(slide, Inches(1.25), y, Inches(2.2), Inches(0.22), item["label"], size=15, bold=True, color=BLUE)
        add_text(slide, Inches(1.25), y + Inches(0.22), Inches(4.55), Inches(0.5), item["text"], size=14, color=GRAY_TEXT)
        y += Inches(1.1)
    right_builder(slide)
    add_slide_number(slide, idx, total)
    add_footer(slide)


def grid_slide(prs, idx, total, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_card(slide, Inches(0.45), Inches(0.45), Inches(12.45), Inches(6.25))
    add_text(slide, Inches(0.95), Inches(0.95), Inches(6.0), Inches(0.5), title, size=24, bold=True, color=BLUE_DARK)
    positions = [
        (Inches(1.0), Inches(1.9)),
        (Inches(6.5), Inches(1.9)),
        (Inches(1.0), Inches(4.1)),
        (Inches(6.5), Inches(4.1)),
    ]
    color_map = [BLUE, BLUE_DARK, GREEN, PURPLE]
    short_map = ["UI", "MAIL", "CLS", "GEN"]
    for i, item in enumerate(items):
        left, top = positions[i]
        add_card(slide, left, top, Inches(5.1), Inches(1.55), fill=BLUE_LIGHT, line=RGBColor(191, 219, 254))
        add_icon_circle(slide, left + Inches(0.28), top + Inches(0.32), Inches(0.56), color_map[i], short_map[i])
        add_text(slide, left + Inches(1.0), top + Inches(0.25), Inches(3.7), Inches(0.3), item["title"], size=18, bold=True, color=BLUE_DARK)
        add_text(slide, left + Inches(1.0), top + Inches(0.72), Inches(3.7), Inches(0.4), item["desc"], size=13, color=GRAY_TEXT)
    add_slide_number(slide, idx, total)
    add_footer(slide)


def diagram_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_card(slide, Inches(0.45), Inches(0.45), Inches(12.45), Inches(6.25))
    add_text(slide, Inches(0.95), Inches(0.95), Inches(6.5), Inches(0.5), "系统总体架构设计", size=24, bold=True, color=BLUE_DARK)
    add_text(slide, Inches(0.95), Inches(1.45), Inches(8.5), Inches(0.3), "系统采用四层架构模型，确保模块解耦与高性能运行：", size=14, color=GRAY_TEXT)

    layers = [
        ("用户界面层 (PyQt5 GUI)", RGBColor(29, 78, 216)),
        ("业务逻辑层 (分类决策/任务调度)", RGBColor(37, 99, 235)),
        ("模型推理层 (Llama 3.2 / Prompt Engineering)", RGBColor(59, 130, 246)),
        ("数据接入层 (IMAP / SMTP / Local DB)", RGBColor(96, 165, 250)),
    ]
    y = Inches(2.0)
    for i, (text, color) in enumerate(layers):
        add_card(slide, Inches(3.15), y, Inches(6.0), Inches(0.55), fill=color, line=color)
        add_text(slide, Inches(3.3), y + Inches(0.08), Inches(5.7), Inches(0.2), text, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(layers) - 1:
            add_text(slide, Inches(6.0), y + Inches(0.62), Inches(0.3), Inches(0.2), "▼", size=18, bold=True, color=RGBColor(147, 197, 253), align=PP_ALIGN.CENTER)
        y += Inches(0.95)
    add_slide_number(slide, idx, total)
    add_footer(slide)


def table_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_card(slide, Inches(0.45), Inches(0.45), Inches(12.45), Inches(6.25))
    add_text(slide, Inches(0.95), Inches(0.95), Inches(6.5), Inches(0.5), "测试结果展示", size=24, bold=True, color=BLUE_DARK)
    headers = ["测试维度", "测试集规模", "准确率/结果", "平均耗时"]
    rows = [
        ["分类任务", "50 封邮件", "96.0%", "0.8s"],
        ["回复生成", "20 轮交互", "语义正确率 100%", "2.5s"],
        ["正文润色", "15 组样本", "通过语法校验", "1.8s"],
        ["系统稳定性", "48h 运行", "无内存溢出/闪退", "-"],
    ]
    table = slide.shapes.add_table(1 + len(rows), len(headers), Inches(0.9), Inches(1.8), Inches(11.4), Inches(3.8)).table
    widths = [2.2, 2.3, 3.8, 1.9]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                fontify(run, size=13, bold=True, color=WHITE)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else GRAY_BG
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    fontify(run, size=12, color=GRAY_TEXT)
    add_slide_number(slide, idx, total)
    add_footer(slide)


def build_right_system_vision(slide):
    add_card(slide, Inches(6.65), Inches(1.8), Inches(5.2), Inches(4.35), fill=BLUE_LIGHT, line=RGBColor(191, 219, 254))
    add_text(slide, Inches(7.55), Inches(2.05), Inches(3.4), Inches(0.3), "系统愿景", size=18, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    items = [
        ("本地私有化部署 (Data Privacy)", GREEN),
        ("低延时自动响应 (Efficiency)", YELLOW),
        ("轻量级模型驱动 (Model Edge)", PURPLE),
    ]
    y = Inches(2.55)
    for text, color in items:
        add_card(slide, Inches(7.25), y, Inches(4.0), Inches(0.7))
        add_icon_circle(slide, Inches(7.5), y + Inches(0.12), Inches(0.42), color, "•")
        add_text(slide, Inches(8.05), y + Inches(0.17), Inches(2.9), Inches(0.22), text, size=11, color=GRAY_TEXT)
        y += Inches(1.0)


def build_right_classifier(slide):
    add_card(slide, Inches(6.65), Inches(1.8), Inches(5.2), Inches(4.35), fill=RGBColor(17, 24, 39), line=RGBColor(17, 24, 39))
    add_text(slide, Inches(7.05), Inches(2.1), Inches(2.2), Inches(0.22), "// Prompt Snippet", size=12, color=GREEN)
    prompt_lines = [
        "\"System: You are an email classifier...\"",
        "\"User: Apply for a discount...\"",
        "\"Output: { \\\"category\\\": \\\"Promotion\\\" }\"",
    ]
    y = Inches(2.6)
    for line in prompt_lines:
        add_text(slide, Inches(7.05), y, Inches(4.1), Inches(0.22), line, size=12, color=RGBColor(74, 222, 128),)
        y += Inches(0.42)


def build_right_summary(slide):
    add_card(slide, Inches(6.65), Inches(1.8), Inches(5.2), Inches(4.35), fill=BLUE_LIGHT, line=RGBColor(191, 219, 254))
    add_text(slide, Inches(7.2), Inches(2.05), Inches(3.9), Inches(0.28), "后续优化方向：", size=17, bold=True, color=BLUE_DARK)
    add_bullets(
        slide,
        Inches(7.15),
        Inches(2.55),
        Inches(4.1),
        Inches(2.5),
        [
            "针对邮件数据集进行 LoRA 微调",
            "增加多语言实时翻译功能",
            "支持多模态附件解析",
        ],
        size=14,
        box=False,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 12
    title_slide(prs, 1, total)
    list_slide(
        prs,
        2,
        total,
        "汇报目录",
        [
            "01 设计目的与意义",
            "02 技术路线",
            "03 系统总体架构设计",
            "04 关键功能模块实现",
            "05 项目创新点",
            "06 测试结果与分析",
            "07 总结与展望",
        ],
    )
    split_slide(
        prs,
        3,
        total,
        "设计目的与意义",
        [
            {"label": "效率挑战", "text": "传统邮件处理依赖人工，耗时较长且易疲劳。"},
            {"label": "隐私保护", "text": "邮件含个人敏感信息，本地化处理可杜绝数据上云风险。"},
            {"label": "轻量应用", "text": "Llama 3.2 1B 使在普通 PC 上运行 LLM 成为可能。"},
        ],
        build_right_system_vision,
    )
    grid_slide(
        prs,
        4,
        total,
        "技术路线",
        [
            {"title": "UI 框架", "desc": "PyQt5 桌面级交互"},
            {"title": "通信协议", "desc": "IMAP/SMTP 邮件接入"},
            {"title": "推理引擎", "desc": "Transformers / Ollama"},
            {"title": "核心模型", "desc": "Llama 3.2 1B (LLM)"},
        ],
    )
    diagram_slide(prs, 5, total)
    grid_slide(
        prs,
        6,
        total,
        "功能模块划分",
        [
            {"title": "账户管理", "desc": "多邮箱接入与协议配置"},
            {"title": "邮件解析", "desc": "HTML/附件/正文提取"},
            {"title": "智能分类", "desc": "规则预筛 + 模型二次判别"},
            {"title": "润色/回复", "desc": "风格迁移与自动草稿"},
        ],
    )
    split_slide(
        prs,
        7,
        total,
        "关键实现：邮件分类逻辑",
        [
            {"label": "混合过滤策略", "text": "第一层：正则表达式/黑名单关键词（秒级过滤典型广告）。"},
            {"label": "模型深感判别", "text": "第二层：LLM 分析语境，识别隐蔽性垃圾邮件。"},
            {"label": "提示词工程", "text": "使用 Few-shot 样本引导模型输出固定格式 JSON。"},
        ],
        build_right_classifier,
    )
    list_slide(
        prs,
        8,
        total,
        "关键实现：智能回复与润色",
        [
            "回复生成：结合邮件原文与用户关键词，自动构建结构化回复模板。",
            "风格迁移：支持自然、正式、商务三种模式。",
            "后处理机制：利用正则和字符串清洗，去除模型生成的冗余语气词。",
            "即时发送：生成预览后支持一键 SMTP 发送。",
        ],
    )
    grid_slide(
        prs,
        9,
        total,
        "项目创新点",
        [
            {"title": "本地化办公", "desc": "完全离线运行，保障隐私"},
            {"title": "极致轻量化", "desc": "在 1B 参数模型上实现生产力闭环"},
            {"title": "双后端支持", "desc": "兼容 Ollama 与 Transformers 环境"},
            {"title": "工程完整性", "desc": "完成从算法到桌面可执行程序的封装"},
        ],
    )
    table_slide(prs, 10, total)
    split_slide(
        prs,
        11,
        total,
        "总结与展望",
        [
            {"label": "已完成工作", "text": "实现了邮件收发、本地模型分类、风格化回复及桌面化封装。"},
            {"label": "研究价值", "text": "验证了轻量化小语言模型在垂类办公场景的实用性。"},
        ],
        build_right_summary,
    )
    title_slide(prs, 12, total)
    last = prs.slides[-1]
    # overwrite final title slide texts to match the html
    title_textboxes = [sh for sh in last.shapes if sh.has_text_frame]
    # easiest: cover central area with white card and rewrite
    add_card(last, Inches(0.95), Inches(1.35), Inches(11.4), Inches(4.2), fill=WHITE, line=WHITE)
    add_icon_circle(last, Inches(6.0), Inches(1.0), Inches(1.1), GREEN, "✓")
    add_text(last, Inches(3.9), Inches(2.2), Inches(5.5), Inches(0.8), "谢 谢 倾 听", size=30, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    add_text(last, Inches(4.25), Inches(3.1), Inches(4.8), Inches(0.45), "敬请各位老师批评指正", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(last, Inches(4.3), Inches(4.3), Inches(4.6), Inches(0.25), "汇报人：高煜同", size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    add_text(last, Inches(4.7), Inches(4.7), Inches(3.7), Inches(0.25), "日期：2026.4.18", size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    try:
        prs.save(OUT)
        print(OUT)
    except PermissionError:
        prs.save(FALLBACK_OUT)
        print(FALLBACK_OUT)


if __name__ == "__main__":
    main()
